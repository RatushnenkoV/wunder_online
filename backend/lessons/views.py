from django.core.exceptions import ValidationError
from django.db import models as django_models
from django.shortcuts import get_object_or_404
from rest_framework.decorators import api_view, permission_classes
from rest_framework.permissions import IsAuthenticated
from rest_framework.response import Response

from django.utils import timezone
from accounts.permissions import PasswordChanged
from core.validators import validate_file_mime, ALLOWED_IMAGES, ALLOWED_PDF, ALLOWED_EXCEL

ALLOWED_PRESENTATION_FILES = ALLOWED_PDF + ['application/zip']  # PDF + PPTX (zip)
from .models import Lesson, LessonFolder, Slide, LessonMedia, LessonSession, FormAnswer, VocabProgress, Textbook, TextbookAnnotation, LessonAssignment
from .serializers import LessonFolderSerializer, LessonSerializer, SlideSerializer, LessonMediaSerializer, LessonSessionSerializer, TextbookSerializer, LessonAssignmentSerializer
from .utils import compute_form_results


def _ctx(request):
    return {'request': request}


def _is_staff(user):
    return user.is_admin or user.is_teacher


def _folder_has_lessons(folder):
    """Рекурсивно проверяет, есть ли уроки в папке или её подпапках."""
    if folder.lessons.exists():
        return True
    for child in folder.children.all():
        if _folder_has_lessons(child):
            return True
    return False


# ─── Папки ────────────────────────────────────────────────────────────────────

@api_view(['GET', 'POST'])
@permission_classes([IsAuthenticated, PasswordChanged])
def folder_list_create(request):
    """GET: мои папки (верхнего уровня). POST: создать папку."""
    if request.method == 'GET':
        folders = LessonFolder.objects.filter(
            owner=request.user, parent=None
        ).prefetch_related('children', 'lessons')
        return Response(LessonFolderSerializer(folders, many=True, context=_ctx(request)).data)

    if not _is_staff(request.user):
        return Response({'error': 'Только учителя могут создавать папки'}, status=403)

    serializer = LessonFolderSerializer(data=request.data, context=_ctx(request))
    if serializer.is_valid():
        serializer.save(owner=request.user)
        return Response(serializer.data, status=201)
    return Response(serializer.errors, status=400)


@api_view(['GET', 'PUT', 'DELETE'])
@permission_classes([IsAuthenticated, PasswordChanged])
def folder_detail(request, folder_id):
    folder = get_object_or_404(LessonFolder, id=folder_id)

    if request.method == 'GET':
        return Response(LessonFolderSerializer(folder, context=_ctx(request)).data)

    # Изменять/удалять может только владелец или admin
    if folder.owner_id != request.user.id and not request.user.is_admin:
        return Response({'error': 'Нет доступа'}, status=403)

    if request.method == 'PUT':
        serializer = LessonFolderSerializer(folder, data=request.data, partial=True, context=_ctx(request))
        if serializer.is_valid():
            serializer.save()
            return Response(serializer.data)
        return Response(serializer.errors, status=400)

    if _folder_has_lessons(folder):
        return Response(
            {'error': 'Нельзя удалить папку, в которой есть уроки. Сначала удалите или переместите все уроки.'},
            status=400,
        )
    folder.delete()
    return Response(status=204)


@api_view(['GET'])
@permission_classes([IsAuthenticated, PasswordChanged])
def folder_contents(request, folder_id):
    """Содержимое папки: вложенные папки + уроки. Доступно всем авторизованным."""
    folder = get_object_or_404(LessonFolder, id=folder_id)

    subfolders = folder.children.all()
    lessons = folder.lessons.all()

    return Response({
        'folder': LessonFolderSerializer(folder, context=_ctx(request)).data,
        'subfolders': LessonFolderSerializer(subfolders, many=True, context=_ctx(request)).data,
        'lessons': LessonSerializer(lessons, many=True, context=_ctx(request)).data,
    })


# ─── Уроки ────────────────────────────────────────────────────────────────────

@api_view(['GET', 'POST'])
@permission_classes([IsAuthenticated, PasswordChanged])
def lesson_list_create(request):
    """
    GET ?tab=mine|all&folder=<id>
    POST: создать урок
    """
    if request.method == 'GET':
        tab = request.query_params.get('tab', 'mine')

        if tab == 'all':
            # Все уроки школы (только для staff/admin)
            if not _is_staff(request.user):
                return Response({'error': 'Нет доступа'}, status=403)
            lessons = Lesson.objects.select_related('owner', 'folder').all()
        else:
            # Только мои уроки
            lessons = Lesson.objects.filter(owner=request.user).select_related('owner', 'folder')

        # ?picker=true — все уроки без фильтра по папке (для выбора в КТП/проектах)
        if request.query_params.get('picker') == 'true':
            return Response(LessonSerializer(lessons.order_by('title'), many=True, context=_ctx(request)).data)

        folder_id = request.query_params.get('folder')
        if folder_id:
            lessons = lessons.filter(folder_id=folder_id)
        else:
            # Только корневые (без папки) если folder не указан
            lessons = lessons.filter(folder=None)

        return Response(LessonSerializer(lessons, many=True, context=_ctx(request)).data)

    # POST — создать урок
    if not _is_staff(request.user):
        return Response({'error': 'Только учителя могут создавать уроки'}, status=403)

    serializer = LessonSerializer(data=request.data, context=_ctx(request))
    if serializer.is_valid():
        lesson = serializer.save(owner=request.user)
        return Response(LessonSerializer(lesson, context=_ctx(request)).data, status=201)
    return Response(serializer.errors, status=400)


@api_view(['GET', 'PUT', 'DELETE'])
@permission_classes([IsAuthenticated, PasswordChanged])
def lesson_detail(request, lesson_id):
    lesson = get_object_or_404(
        Lesson.objects.select_related('owner', 'folder'),
        id=lesson_id,
    )

    if request.method == 'GET':
        return Response(LessonSerializer(lesson, context=_ctx(request)).data)

    # Изменять/удалять может владелец или admin
    if lesson.owner_id != request.user.id and not request.user.is_admin:
        return Response({'error': 'Нет доступа'}, status=403)

    if request.method == 'PUT':
        serializer = LessonSerializer(lesson, data=request.data, partial=True, context=_ctx(request))
        if serializer.is_valid():
            lesson = serializer.save()
            return Response(LessonSerializer(lesson, context=_ctx(request)).data)
        return Response(serializer.errors, status=400)

    lesson.delete()
    return Response(status=204)


# ─── Импорт презентаций ───────────────────────────────────────────────────────

_IMPORT_W = 960
_IMPORT_H = 540
# Имена тем, которые не несут реального имени шрифта
_THEME_FONT_PLACEHOLDERS = {'+mj-lt', '+mn-lt', '+mj-ea', '+mn-ea', '+mj-cs', '+mn-cs'}


def _qn(tag: str) -> str:
    """Clark-нотация для OOXML-тегов, например 'a:rPr' → '{...ns}rPr'."""
    from pptx.oxml.ns import qn
    return qn(tag)


# ── XML-помощники для чтения унаследованных свойств ──────────────────────────

def _xml_sz(rPr, pPr, txBody) -> float | None:
    """Размер шрифта в pt: rPr.sz → pPr.defRPr.sz → lstStyle.lvl1pPr.*"""
    def _sz_from(el):
        if el is None:
            return None
        sz = el.get('sz')
        return int(sz) / 100 if sz else None

    v = _sz_from(rPr)
    if v:
        return v
    if pPr is not None:
        v = _sz_from(pPr.find(_qn('a:defRPr')))
        if v:
            return v
    if txBody is not None:
        lst = txBody.find(_qn('a:lstStyle'))
        if lst is not None:
            for tag in ('a:lvl1pPr', 'a:lvl2pPr'):
                lvl = lst.find(_qn(tag))
                if lvl is not None:
                    v = _sz_from(lvl.find(_qn('a:defRPr'))) or _sz_from(lvl)
                    if v:
                        return v
    return None


def _xml_font_name(rPr, pPr, txBody) -> str | None:
    """Имя шрифта из цепочки rPr → pPr.defRPr → lstStyle."""
    def _name_from(el):
        if el is None:
            return None
        lat = el.find(_qn('a:latin'))
        if lat is not None:
            tf = lat.get('typeface')
            if tf and tf not in _THEME_FONT_PLACEHOLDERS:
                return tf
        return None

    v = _name_from(rPr)
    if v:
        return v
    if pPr is not None:
        v = _name_from(pPr.find(_qn('a:defRPr')))
        if v:
            return v
    if txBody is not None:
        lst = txBody.find(_qn('a:lstStyle'))
        if lst is not None:
            lvl = lst.find(_qn('a:lvl1pPr'))
            if lvl is not None:
                v = _name_from(lvl.find(_qn('a:defRPr')))
                if v:
                    return v
    return None


def _xml_color_hex(rPr, pPr) -> str | None:
    """Цвет из solidFill/srgbClr: rPr → pPr.defRPr."""
    def _from(el):
        if el is None:
            return None
        sf = el.find(_qn('a:solidFill'))
        if sf is not None:
            sc = sf.find(_qn('a:srgbClr'))
            if sc is not None:
                val = sc.get('val')
                if val:
                    return f'#{val}'
        return None

    v = _from(rPr)
    if v:
        return v
    if pPr is not None:
        v = _from(pPr.find(_qn('a:defRPr')))
        if v:
            return v
    return None


def _xml_bool(el, attr: str) -> bool | None:
    """Читает булев атрибут OOXML ('0'/'1'/'true'/'false') из элемента."""
    if el is None:
        return None
    val = el.get(attr)
    if val is None:
        return None
    return val not in ('0', 'false')


# ── Конвертация параграфов ────────────────────────────────────────────────────

def _run_to_html(run, para, txBody) -> str:
    """Фрагмент PPTX → HTML <span> с полными инлайн-стилями (через XML)."""
    text = run.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
    if not text:
        return ''

    rPr = run._r.find(_qn('a:rPr'))
    pPr = para._p.find(_qn('a:pPr'))
    styles = []

    # Шрифт
    name = _xml_font_name(rPr, pPr, txBody)
    if name:
        safe_name = name.replace("'", "\\'")
        styles.append(f"font-family: '{safe_name}', sans-serif")

    # Размер (pt ≈ px на стандартном холсте 960px/10")
    sz = _xml_sz(rPr, pPr, txBody)
    if sz and sz > 0:
        styles.append(f'font-size: {round(sz)}px')

    # Цвет
    color = _xml_color_hex(rPr, pPr)
    if color:
        styles.append(f'color: {color}')

    # Начертание
    if _xml_bool(rPr, 'b'):
        styles.append('font-weight: bold')
    if _xml_bool(rPr, 'i'):
        styles.append('font-style: italic')
    u = rPr.get('u') if rPr is not None else None
    if u and u != 'none':
        styles.append('text-decoration: underline')

    if styles:
        return f'<span style="{"; ".join(styles)}">{text}</span>'
    return text


def _para_to_html(para, txBody) -> str:
    """Параграф PPTX → HTML <p> с выравниванием и форматированием."""
    try:
        from pptx.enum.text import PP_ALIGN
    except Exception:
        PP_ALIGN = None

    inner = ''.join(_run_to_html(r, para, txBody) for r in para.runs)
    if not inner:
        raw = para.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        if not raw.strip():
            return ''
        inner = raw

    align_style = ''
    if PP_ALIGN:
        try:
            a = para.alignment
            if a == PP_ALIGN.CENTER:
                align_style = 'text-align: center'
            elif a == PP_ALIGN.RIGHT:
                align_style = 'text-align: right'
            elif a == PP_ALIGN.JUSTIFY:
                align_style = 'text-align: justify'
        except Exception:
            pass

    attrs = f' style="{align_style}"' if align_style else ''
    return f'<p{attrs}>{inner}</p>'


# ── Фон слайда ────────────────────────────────────────────────────────────────

def _get_theme_scheme(pptx_slide):
    """Возвращает элемент <a:clrScheme> из темы мастера слайда."""
    try:
        from lxml import etree
        master_part = pptx_slide.slide_layout.slide_master.part
        theme_rtype = (
            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme'
        )
        for rel in master_part.rels.values():
            if rel.reltype == theme_rtype:
                blob = rel._target.blob
                tree = etree.fromstring(blob)
                return tree.find('.//' + _qn('a:clrScheme'))
    except Exception:
        pass
    return None


def _resolve_scheme_color(scheme_el, val_name: str) -> str | None:
    """Разрешает имя schemeClr (bg1/acc1/dk1/…) в HEX-строку через <a:clrScheme>."""
    name_map = {
        'dk1': 'a:dk1',    'lt1': 'a:lt1',
        'dk2': 'a:dk2',    'lt2': 'a:lt2',
        'acc1': 'a:accent1', 'acc2': 'a:accent2',
        'acc3': 'a:accent3', 'acc4': 'a:accent4',
        'acc5': 'a:accent5', 'acc6': 'a:accent6',
        'hlink': 'a:hlink', 'folHlink': 'a:folHlink',
        # псевдонимы
        'bg1': 'a:lt1', 'bg2': 'a:lt2',
        'tx1': 'a:dk1', 'tx2': 'a:dk2',
        'accent1': 'a:accent1', 'accent2': 'a:accent2',
        'accent3': 'a:accent3', 'accent4': 'a:accent4',
        'accent5': 'a:accent5', 'accent6': 'a:accent6',
    }
    tag = name_map.get(val_name)
    if not tag or scheme_el is None:
        return None
    slot = scheme_el.find(_qn(tag))
    if slot is None:
        return None
    rgb_el = slot.find(_qn('a:srgbClr'))
    if rgb_el is not None:
        return rgb_el.get('val')
    sys_el = slot.find(_qn('a:sysClr'))
    if sys_el is not None:
        return sys_el.get('lastClr')
    return None


def _apply_color_mods(hex_rgb: str, mods_el) -> str:
    """Применяет OOXML-модификаторы цвета (lumMod/lumOff/shade/tint) к HEX-строке."""
    import colorsys
    try:
        r = int(hex_rgb[0:2], 16) / 255.0
        g = int(hex_rgb[2:4], 16) / 255.0
        b = int(hex_rgb[4:6], 16) / 255.0

        def pct(el, attr):
            v = el.get(attr)
            return int(v) / 100000.0 if v else None

        for child in mods_el:
            local = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if local == 'lumMod':
                f = pct(child, 'val')
                if f is not None:
                    h, l, s = colorsys.rgb_to_hls(r, g, b)
                    l = max(0.0, min(1.0, l * f))
                    r, g, b = colorsys.hls_to_rgb(h, l, s)
            elif local == 'lumOff':
                f = pct(child, 'val')
                if f is not None:
                    h, l, s = colorsys.rgb_to_hls(r, g, b)
                    l = max(0.0, min(1.0, l + f))
                    r, g, b = colorsys.hls_to_rgb(h, l, s)
            elif local == 'shade':
                f = pct(child, 'val')
                if f is not None:
                    r = max(0.0, min(1.0, r * f))
                    g = max(0.0, min(1.0, g * f))
                    b = max(0.0, min(1.0, b * f))
            elif local == 'tint':
                f = pct(child, 'val')
                if f is not None:
                    r = max(0.0, min(1.0, r * f + (1.0 - f)))
                    g = max(0.0, min(1.0, g * f + (1.0 - f)))
                    b = max(0.0, min(1.0, b * f + (1.0 - f)))

        return f'{int(round(r * 255)):02X}{int(round(g * 255)):02X}{int(round(b * 255)):02X}'
    except Exception:
        return hex_rgb


def _color_from_fill_parent(el, scheme_el) -> str | None:
    """Извлекает цвет из элемента, содержащего solidFill/gradFill (например, bgPr или bgRef)."""
    # solidFill
    sf = el.find(_qn('a:solidFill'))
    if sf is not None:
        srgb = sf.find(_qn('a:srgbClr'))
        if srgb is not None and srgb.get('val'):
            return f'#{srgb.get("val")}'
        sch = sf.find(_qn('a:schemeClr'))
        if sch is not None:
            base = _resolve_scheme_color(scheme_el, sch.get('val', ''))
            if base:
                return f'#{_apply_color_mods(base, sch)}'
        sys_el = sf.find(_qn('a:sysClr'))
        if sys_el is not None and sys_el.get('lastClr'):
            return f'#{sys_el.get("lastClr")}'

    # gradFill: берём первую остановку
    gf = el.find(_qn('a:gradFill'))
    if gf is not None:
        lst = gf.find(_qn('a:gsLst'))
        if lst is not None:
            stops = list(lst)
            if stops:
                first = stops[0]
                srgb = first.find(_qn('a:srgbClr'))
                if srgb is not None and srgb.get('val'):
                    return f'#{srgb.get("val")}'
                sch = first.find(_qn('a:schemeClr'))
                if sch is not None:
                    base = _resolve_scheme_color(scheme_el, sch.get('val', ''))
                    if base:
                        return f'#{_apply_color_mods(base, sch)}'
                sys_el = first.find(_qn('a:sysClr'))
                if sys_el is not None and sys_el.get('lastClr'):
                    return f'#{sys_el.get("lastClr")}'

    # schemeClr / srgbClr непосредственно в el (для bgRef)
    sch = el.find(_qn('a:schemeClr'))
    if sch is not None:
        base = _resolve_scheme_color(scheme_el, sch.get('val', ''))
        if base:
            return f'#{_apply_color_mods(base, sch)}'
    srgb = el.find(_qn('a:srgbClr'))
    if srgb is not None and srgb.get('val'):
        return f'#{srgb.get("val")}'

    return None


def _bg_from_csld(cSld_el, scheme_el) -> str | None:
    """Извлекает цвет фона из элемента <p:cSld>."""
    if cSld_el is None:
        return None
    bg = cSld_el.find(_qn('p:bg'))
    if bg is None:
        return None
    bgPr = bg.find(_qn('p:bgPr'))
    if bgPr is not None:
        color = _color_from_fill_parent(bgPr, scheme_el)
        if color:
            return color
    bgRef = bg.find(_qn('p:bgRef'))
    if bgRef is not None:
        return _color_from_fill_parent(bgRef, scheme_el)
    return None


def _slide_bg_color(pptx_slide) -> str | None:
    """Цвет фона слайда: slide → layout → master, поддержка schemeClr/srgbClr/gradFill."""
    try:
        scheme_el = _get_theme_scheme(pptx_slide)

        # 1. Слайд
        cSld = pptx_slide._element.find(_qn('p:cSld'))
        color = _bg_from_csld(cSld, scheme_el)
        if color:
            return color

        # 2. Макет (layout)
        try:
            layout_cSld = pptx_slide.slide_layout._element.find(_qn('p:cSld'))
            color = _bg_from_csld(layout_cSld, scheme_el)
            if color:
                return color
        except Exception:
            pass

        # 3. Мастер — возвращаем только нестандартные цвета (не белый)
        try:
            master_cSld = pptx_slide.slide_layout.slide_master._element.find(_qn('p:cSld'))
            color = _bg_from_csld(master_cSld, scheme_el)
            if color and color.upper() not in ('#FFFFFF', '#FFF'):
                return color
        except Exception:
            pass
    except Exception:
        pass
    return None


# ── Цвет/тип авто-фигур ───────────────────────────────────────────────────────

def _shape_fill_stroke(shape):
    """Возвращает (fill_color, stroke_color, stroke_width_px)."""
    fill_color = 'transparent'
    stroke_color = 'transparent'
    stroke_width = 3

    try:
        ft = shape.fill.type
        if ft is not None:
            try:
                fill_color = f'#{shape.fill.fore_color.rgb}'
            except Exception:
                fill_color = '#6366f1'
    except Exception:
        pass

    try:
        line = shape.line
        try:
            if line.color.type is not None:
                stroke_color = f'#{line.color.rgb}'
        except Exception:
            pass
        if stroke_color == 'transparent':
            try:
                if line.width and line.width > 0:
                    stroke_color = '#374151'
            except Exception:
                pass
        try:
            if line.width:
                stroke_width = max(1, round(line.width / 12700))
        except Exception:
            pass
    except Exception:
        pass

    return fill_color, stroke_color, stroke_width


def _canvas_shape(shape) -> str:
    """Тип нашего canvas-блока: rect/circle/triangle/diamond/star/line."""
    try:
        s = str(shape.auto_shape_type).upper()
        if any(k in s for k in ('OVAL', 'ELLIPSE', 'CIRCLE')):
            return 'circle'
        if any(k in s for k in ('TRIANGLE', 'ISOSCELES', 'RIGHT_TRIANGLE')):
            return 'triangle'
        if 'DIAMOND' in s:
            return 'diamond'
        if 'STAR' in s:
            return 'star'
    except Exception:
        pass
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            return 'line'
    except Exception:
        pass
    return 'rect'


# ── PDF ───────────────────────────────────────────────────────────────────────

def _import_pdf(request, lesson, file_obj):
    """Импорт PDF: каждая страница → content-слайд с блоком-изображением на весь холст."""
    import fitz  # pymupdf
    from django.core.files.base import ContentFile

    data = file_obj.read()
    doc = fitz.open(stream=data, filetype='pdf')
    try:
        for i, page in enumerate(doc):
            mat = fitz.Matrix(2, 2)  # 2× для чёткости
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes('png')

            media = LessonMedia(lesson=lesson)
            media.file.save(f'page_{i + 1}.png', ContentFile(img_bytes), save=True)
            img_url = request.build_absolute_uri(media.file.url)

            Slide.objects.create(
                lesson=lesson,
                order=i,
                slide_type=Slide.TYPE_CONTENT,
                content={
                    'blocks': [{
                        'id': f'b{i}_1',
                        'type': 'image',
                        'x': 0, 'y': 0,
                        'w': _IMPORT_W, 'h': _IMPORT_H,
                        'zIndex': 1,
                        'rotation': 0,
                        'src': img_url,
                        'alt': f'Страница {i + 1}',
                    }],
                },
            )
    finally:
        doc.close()


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _import_pptx(request, lesson, file_obj):
    """
    Импорт PPTX: каждый слайд → content-слайд.
    Поддерживает: текст (шрифт/размер/цвет/начертание/выравнивание),
    изображения, авто-фигуры (rect/circle/triangle/diamond/star/line),
    цвет фона слайда, угол поворота объектов.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from django.core.files.base import ContentFile

    AUTO_SHAPE_TYPES = (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM)

    prs = Presentation(file_obj)
    slide_w = prs.slide_width or 1
    slide_h = prs.slide_height or 1

    for i, pptx_slide in enumerate(prs.slides):
        blocks = []
        idx = 1
        bg_color = _slide_bg_color(pptx_slide)

        for shape in pptx_slide.shapes:
            try:
                x = int(shape.left / slide_w * _IMPORT_W)
                y = int(shape.top  / slide_h * _IMPORT_H)
                w = max(50, int(shape.width  / slide_w * _IMPORT_W))
                h = max(20, int(shape.height / slide_h * _IMPORT_H))
                rotation = float(getattr(shape, 'rotation', 0) or 0)
                stype = getattr(shape, 'shape_type', None)

                # ── Изображение ──────────────────────────────────────────────
                if stype == MSO_SHAPE_TYPE.PICTURE:
                    img_data = shape.image.blob
                    img_ext = (shape.image.ext or 'png').lstrip('.')
                    media = LessonMedia(lesson=lesson)
                    media.file.save(f'slide_{i}_{idx}.{img_ext}', ContentFile(img_data), save=True)
                    blocks.append({
                        'id': f'b{i}_{idx}', 'type': 'image',
                        'x': max(0, x), 'y': max(0, y),
                        'w': min(w, _IMPORT_W), 'h': min(h, _IMPORT_H),
                        'zIndex': idx, 'rotation': rotation,
                        'src': request.build_absolute_uri(media.file.url), 'alt': '',
                    })
                    idx += 1
                    continue

                # ── Авто-фигура (без текста или с заливкой) ──────────────────
                if stype in AUTO_SHAPE_TYPES or stype == MSO_SHAPE_TYPE.LINE:
                    fc, sc, sw = _shape_fill_stroke(shape)
                    if fc != 'transparent' or sc != 'transparent':
                        blocks.append({
                            'id': f'b{i}_{idx}', 'type': 'shape',
                            'shape': _canvas_shape(shape),
                            'x': max(0, x), 'y': max(0, y),
                            'w': min(w, _IMPORT_W), 'h': min(h, _IMPORT_H),
                            'zIndex': idx, 'rotation': rotation,
                            'fillColor': fc, 'strokeColor': sc, 'strokeWidth': sw,
                        })
                        idx += 1

                # ── Текст (присутствует на любом типе фигуры) ─────────────────
                if shape.has_text_frame:
                    txBody = shape.text_frame._txBody
                    parts = [_para_to_html(p, txBody) for p in shape.text_frame.paragraphs]
                    html = ''.join(p for p in parts if p)
                    if html:
                        blocks.append({
                            'id': f'b{i}_{idx}', 'type': 'text',
                            'x': max(0, x), 'y': max(0, y),
                            'w': min(w, _IMPORT_W), 'h': min(h, _IMPORT_H),
                            'zIndex': idx, 'rotation': rotation,
                            'html': html,
                        })
                        idx += 1

            except Exception:
                continue

        content: dict = {'blocks': blocks}
        if bg_color:
            content['background'] = bg_color

        Slide.objects.create(
            lesson=lesson,
            order=i,
            slide_type=Slide.TYPE_CONTENT,
            content=content,
        )


@api_view(['POST'])
@permission_classes([IsAuthenticated, PasswordChanged])
def import_presentation(request):
    """POST /lessons/import/ — создать урок из файла PDF или PPTX."""
    if not _is_staff(request.user):
        return Response({'error': 'Только учителя могут создавать уроки'}, status=403)

    file = request.FILES.get('file')
    if not file:
        return Response({'error': 'Файл не указан'}, status=400)

    ext = file.name.rsplit('.', 1)[-1].lower() if '.' in file.name else ''
    if ext not in ('pdf', 'pptx', 'ppt'):
        return Response({'error': 'Поддерживаются только файлы PDF и PPTX'}, status=400)

    try:
        validate_file_mime(file, ALLOWED_PRESENTATION_FILES, label='файл презентации')
    except ValidationError as e:
        return Response({'error': str(e)}, status=400)

    title = (request.data.get('title') or file.name.rsplit('.', 1)[0])[:300]
    folder_id = request.data.get('folder') or None
    cover_color = request.data.get('cover_color', '#6366f1')

    lesson = Lesson.objects.create(
        title=title,
        owner=request.user,
        folder_id=folder_id,
        cover_color=cover_color,
    )

    try:
        if ext == 'pdf':
            _import_pdf(request, lesson, file)
        else:
            _import_pptx(request, lesson, file)
    except Exception as e:
        lesson.delete()
        return Response({'error': f'Ошибка импорта: {str(e)}'}, status=500)

    return Response(LessonSerializer(lesson, context=_ctx(request)).data, status=201)


# ─── Слайды ───────────────────────────────────────────────────────────────────

def _can_edit_lesson(lesson, user):
    return lesson.owner_id == user.id or user.is_admin


@api_view(['GET', 'POST'])
@permission_classes([IsAuthenticated, PasswordChanged])
def slide_list_create(request, lesson_id):
    lesson = get_object_or_404(Lesson.objects.select_related('owner'), id=lesson_id)

    if request.method == 'GET':
        slides = lesson.slides.all()
        return Response(SlideSerializer(slides, many=True, context=_ctx(request)).data)

    if not _can_edit_lesson(lesson, request.user):
        return Response({'error': 'Нет доступа'}, status=403)

    max_order = lesson.slides.aggregate(m=django_models.Max('order'))['m']
    order = (max_order or 0) + 1

    slide = Slide.objects.create(
        lesson=lesson,
        order=order,
        slide_type=request.data.get('slide_type', Slide.TYPE_CONTENT),
        title=request.data.get('title', ''),
        content=request.data.get('content', {}),
    )
    return Response(SlideSerializer(slide, context=_ctx(request)).data, status=201)


@api_view(['GET', 'PUT', 'DELETE'])
@permission_classes([IsAuthenticated, PasswordChanged])
def slide_detail(request, lesson_id, slide_id):
    lesson = get_object_or_404(Lesson, id=lesson_id)
    slide = get_object_or_404(Slide, id=slide_id, lesson=lesson)

    if request.method == 'GET':
        return Response(SlideSerializer(slide, context=_ctx(request)).data)

    if not _can_edit_lesson(lesson, request.user):
        return Response({'error': 'Нет доступа'}, status=403)

    if request.method == 'PUT':
        if 'title' in request.data:
            slide.title = request.data['title']
        if 'slide_type' in request.data:
            slide.slide_type = request.data['slide_type']
        if 'content' in request.data:
            slide.content = request.data['content']
        slide.save()
        return Response(SlideSerializer(slide, context=_ctx(request)).data)

    slide.delete()
    return Response(status=204)


@api_view(['POST'])
@permission_classes([IsAuthenticated, PasswordChanged])
def slides_reorder(request, lesson_id):
    lesson = get_object_or_404(Lesson, id=lesson_id)
    if not _can_edit_lesson(lesson, request.user):
        return Response({'error': 'Нет доступа'}, status=403)

    order_ids = request.data.get('order', [])
    for idx, sid in enumerate(order_ids):
        Slide.objects.filter(id=sid, lesson=lesson).update(order=idx)

    slides = lesson.slides.all()
    return Response(SlideSerializer(slides, many=True, context=_ctx(request)).data)


@api_view(['POST'])
@permission_classes([IsAuthenticated, PasswordChanged])
def slide_image_upload(request, lesson_id, slide_id):
    lesson = get_object_or_404(Lesson, id=lesson_id)
    slide = get_object_or_404(Slide, id=slide_id, lesson=lesson)

    if not _can_edit_lesson(lesson, request.user):
        return Response({'error': 'Нет доступа'}, status=403)

    img = request.FILES.get('image')
    if not img:
        return Response({'error': 'Файл не передан'}, status=400)

    try:
        validate_file_mime(img, ALLOWED_IMAGES, label='изображение слайда')
    except ValidationError as e:
        return Response({'error': str(e)}, status=400)

    if slide.image:
        slide.image.delete(save=False)
    slide.image = img
    slide.save()
    return Response(SlideSerializer(slide, context=_ctx(request)).data)


# ─── Медиафайлы ───────────────────────────────────────────────────────────────

@api_view(['POST'])
@permission_classes([IsAuthenticated, PasswordChanged])
def upload_media(request, lesson_id):
    """Загрузить медиафайл (изображение) для блока на слайде."""
    lesson = get_object_or_404(Lesson, id=lesson_id)
    if not _can_edit_lesson(lesson, request.user):
        return Response({'error': 'Нет доступа'}, status=403)

    f = request.FILES.get('file')
    if not f:
        return Response({'error': 'Файл не передан'}, status=400)

    try:
        validate_file_mime(f, ALLOWED_IMAGES, label='медиафайл урока')
    except ValidationError as e:
        return Response({'error': str(e)}, status=400)

    media = LessonMedia.objects.create(lesson=lesson, file=f)
    return Response(LessonMediaSerializer(media, context=_ctx(request)).data, status=201)


# ─── Дублирование ─────────────────────────────────────────────────────────────

@api_view(['POST'])
@permission_classes([IsAuthenticated, PasswordChanged])
def lesson_duplicate(request, lesson_id):
    """Дублировать урок."""
    lesson = get_object_or_404(Lesson, id=lesson_id)

    if lesson.owner_id != request.user.id and not request.user.is_admin:
        return Response({'error': 'Нет доступа'}, status=403)

    new_lesson = Lesson.objects.create(
        title=f'{lesson.title} (копия)',
        description=lesson.description,
        owner=request.user,
        folder=lesson.folder,
        is_public=False,
        cover_color=lesson.cover_color,
    )
    return Response(LessonSerializer(new_lesson, context=_ctx(request)).data, status=201)


# ─── Сессии ───────────────────────────────────────────────────────────────────

@api_view(['GET', 'POST'])
@permission_classes([IsAuthenticated, PasswordChanged])
def session_list_create(request):
    """
    GET  – список сессий (фильтр ?lesson=id для учителя).
    POST – создать сессию (только учитель/admin).
    """
    if request.method == 'GET':
        lesson_id = request.query_params.get('lesson')
        if lesson_id:
            sessions = LessonSession.objects.filter(lesson_id=lesson_id)
            if not _is_staff(request.user):
                sessions = sessions.filter(is_active=True)
        else:
            if _is_staff(request.user):
                sessions = LessonSession.objects.filter(teacher=request.user)
            else:
                sessions = LessonSession.objects.none()
        return Response(LessonSessionSerializer(sessions, many=True, context=_ctx(request)).data)

    # POST
    if not _is_staff(request.user):
        return Response({'error': 'Только учителя могут начинать уроки'}, status=403)

    lesson_id = request.data.get('lesson')
    school_class_id = request.data.get('school_class')

    lesson = get_object_or_404(Lesson, id=lesson_id)
    first_slide = lesson.slides.first()

    session = LessonSession.objects.create(
        lesson=lesson,
        teacher=request.user,
        school_class_id=school_class_id,
        current_slide=first_slide,
        is_active=True,
    )
    return Response(LessonSessionSerializer(session, context=_ctx(request)).data, status=201)


@api_view(['GET', 'PATCH', 'DELETE'])
@permission_classes([IsAuthenticated, PasswordChanged])
def session_detail(request, session_id):
    session = get_object_or_404(LessonSession, id=session_id)

    if request.method == 'GET':
        return Response(LessonSessionSerializer(session, context=_ctx(request)).data)

    # Изменять/удалять может только учитель сессии или admin
    if session.teacher_id != request.user.id and not request.user.is_admin:
        return Response({'error': 'Нет доступа'}, status=403)

    if request.method == 'PATCH':
        if 'is_active' in request.data and not request.data['is_active']:
            session.is_active = False
            session.ended_at = timezone.now()
        if 'current_slide' in request.data:
            session.current_slide_id = request.data['current_slide']
        session.save()
        return Response(LessonSessionSerializer(session, context=_ctx(request)).data)

    if request.method == 'DELETE':
        session.delete()
        return Response(status=204)


@api_view(['GET'])
@permission_classes([IsAuthenticated, PasswordChanged])
def session_form_results(request, session_id, slide_id):
    """GET результаты формы для учителя (первоначальная загрузка)."""
    session = get_object_or_404(LessonSession, id=session_id)
    if not _is_staff(request.user):
        return Response({'error': 'Нет доступа'}, status=403)
    slide = get_object_or_404(Slide, id=slide_id)
    fa_qs = FormAnswer.objects.filter(session=session, slide=slide).select_related('student')
    return Response(compute_form_results(slide, fa_qs))


@api_view(['GET'])
@permission_classes([IsAuthenticated, PasswordChanged])
def sessions_active(request):
    """Активные сессии, доступные текущему пользователю."""
    user = request.user

    if _is_staff(user):
        # Учителя/адмики видят все активные сессии
        sessions = LessonSession.objects.filter(is_active=True).select_related('lesson', 'teacher', 'school_class')
    elif user.is_student:
        # Студенты видят сессии своего класса
        try:
            from school.models import StudentProfile
            profile = StudentProfile.objects.get(user=user)
            sessions = LessonSession.objects.filter(
                is_active=True,
                school_class=profile.school_class,
            ).select_related('lesson', 'teacher', 'school_class')
        except Exception:
            sessions = LessonSession.objects.none()
    else:
        sessions = LessonSession.objects.none()

    return Response(LessonSessionSerializer(sessions, many=True, context=_ctx(request)).data)


# ─── Словарный прогресс ────────────────────────────────────────────────────────

@api_view(['GET', 'POST'])
@permission_classes([IsAuthenticated, PasswordChanged])
def vocab_progress(request, session_id, slide_id):
    """
    GET  — прогресс всего класса (только для учителей/admin).
    POST — сохранить/обновить прогресс ученика по слову.
    """
    session = get_object_or_404(LessonSession, id=session_id)
    slide = get_object_or_404(Slide, id=slide_id)

    if request.method == 'GET':
        if not _is_staff(request.user):
            return Response({'error': 'Нет доступа'}, status=403)
        rows = VocabProgress.objects.filter(
            session=session, slide=slide,
        ).select_related('student')
        data = [
            {
                'student_id': r.student_id,
                'student_name': f'{r.student.first_name} {r.student.last_name}'.strip(),
                'word_id': r.word_id,
                'attempts': r.attempts,
                'correct': r.correct,
                'learned': r.learned,
                'updated_at': r.updated_at,
            }
            for r in rows
        ]
        return Response(data)

    # POST — ученик сохраняет прогресс
    word_id = request.data.get('word_id')
    if not word_id:
        return Response({'error': 'word_id обязателен'}, status=400)

    attempts = int(request.data.get('attempts', 0))
    correct = int(request.data.get('correct', 0))
    learned = bool(request.data.get('learned', False))

    obj, _ = VocabProgress.objects.update_or_create(
        session=session,
        slide=slide,
        student=request.user,
        word_id=word_id,
        defaults={'attempts': attempts, 'correct': correct, 'learned': learned},
    )
    return Response({
        'student_id': obj.student_id,
        'word_id': obj.word_id,
        'attempts': obj.attempts,
        'correct': obj.correct,
        'learned': obj.learned,
        'updated_at': obj.updated_at,
    }, status=201)


# ─── Обзор уроков школы (вкладка «Все уроки») ─────────────────────────────────

@api_view(['GET'])
@permission_classes([IsAuthenticated, PasswordChanged])
def school_lessons_overview(request):
    """Список пользователей, у которых есть уроки или папки."""
    from accounts.models import User

    teacher_ids = set(
        list(Lesson.objects.values_list('owner_id', flat=True).distinct()) +
        list(LessonFolder.objects.values_list('owner_id', flat=True).distinct())
    )
    users = User.objects.filter(id__in=teacher_ids).order_by('last_name', 'first_name')

    result = []
    for u in users:
        result.append({
            'teacher_id': u.id,
            'teacher_name': f'{u.last_name} {u.first_name}'.strip(),
            'folders_count': LessonFolder.objects.filter(owner=u, parent=None).count(),
            'lessons_count': Lesson.objects.filter(owner=u).count(),
        })
    return Response(result)


@api_view(['GET'])
@permission_classes([IsAuthenticated, PasswordChanged])
def teacher_root_content(request):
    """Корневые папки и уроки конкретного пользователя."""
    from accounts.models import User

    teacher_id = request.query_params.get('teacher_id')
    if not teacher_id:
        return Response({'error': 'teacher_id обязателен'}, status=400)

    teacher = get_object_or_404(User, id=teacher_id)
    folders = LessonFolder.objects.filter(owner=teacher, parent=None)
    lessons = Lesson.objects.filter(owner=teacher, folder=None).select_related('owner')

    return Response({
        'teacher_id': teacher.id,
        'teacher_name': f'{teacher.last_name} {teacher.first_name}'.strip(),
        'folders': LessonFolderSerializer(folders, many=True, context=_ctx(request)).data,
        'lessons': LessonSerializer(lessons, many=True, context=_ctx(request)).data,
    })


# ─── Учебники ─────────────────────────────────────────────────────────────────

@api_view(['GET', 'POST'])
@permission_classes([IsAuthenticated, PasswordChanged])
def textbook_list_create(request):
    """
    GET ?grade_level_id=<id> — список учебников (фильтр по параллели).
    POST — загрузить учебник (только staff).
    """
    if request.method == 'GET':
        qs = Textbook.objects.prefetch_related('grade_levels', 'subject').select_related('uploaded_by')
        grade_level_id = request.query_params.get('grade_level_id')
        if grade_level_id:
            qs = qs.filter(grade_levels__id=grade_level_id)
        return Response(TextbookSerializer(qs, many=True, context=_ctx(request)).data)

    if not _is_staff(request.user):
        return Response({'error': 'Только учителя могут загружать учебники'}, status=403)

    file = request.FILES.get('file')
    if not file:
        return Response({'error': 'Файл обязателен'}, status=400)

    try:
        validate_file_mime(file, ALLOWED_PDF, label='учебник')
    except ValidationError as e:
        return Response({'error': str(e)}, status=400)

    title = request.data.get('title', '').strip() or file.name
    subject_id = request.data.get('subject') or None
    grade_level_ids = request.data.getlist('grade_level_ids')

    textbook = Textbook.objects.create(
        title=title,
        file=file,
        original_name=file.name,
        file_size=file.size,
        subject_id=subject_id,
        uploaded_by=request.user,
    )
    if grade_level_ids:
        textbook.grade_levels.set([int(g) for g in grade_level_ids])

    return Response(TextbookSerializer(textbook, context=_ctx(request)).data, status=201)


@api_view(['GET', 'PUT', 'DELETE'])
@permission_classes([IsAuthenticated, PasswordChanged])
def textbook_detail(request, textbook_id):
    textbook = get_object_or_404(
        Textbook.objects.prefetch_related('grade_levels').select_related('subject', 'uploaded_by'),
        id=textbook_id,
    )

    if request.method == 'GET':
        return Response(TextbookSerializer(textbook, context=_ctx(request)).data)

    if not _is_staff(request.user):
        return Response({'error': 'Нет доступа'}, status=403)

    if request.method == 'PUT':
        title = request.data.get('title', textbook.title).strip() or textbook.title
        subject_id = request.data.get('subject') or None
        grade_level_ids = request.data.getlist('grade_level_ids')

        textbook.title = title
        textbook.subject_id = subject_id
        textbook.save()
        if grade_level_ids is not None:
            textbook.grade_levels.set([int(g) for g in grade_level_ids])

        return Response(TextbookSerializer(textbook, context=_ctx(request)).data)

    textbook.delete()
    return Response(status=204)


@api_view(['GET'])
@permission_classes([IsAuthenticated, PasswordChanged])
def textbook_grade_levels(request):
    """
    Список параллелей, для которых пользователь может видеть учебники.
    Staff: все параллели. Ученик: своя параллель. Родитель: параллели детей.
    """
    from school.models import GradeLevel

    user = request.user

    if user.is_admin or user.is_teacher:
        grade_levels = GradeLevel.objects.order_by('number')
    elif user.is_student:
        try:
            sp = user.student_profile
            gl_id = sp.school_class.grade_level_id if sp.school_class else None
            grade_levels = GradeLevel.objects.filter(id=gl_id) if gl_id else GradeLevel.objects.none()
        except Exception:
            grade_levels = GradeLevel.objects.none()
    elif user.is_parent:
        try:
            children = user.parent_profile.children.select_related('school_class__grade_level')
            gl_ids = set(
                c.school_class.grade_level_id for c in children
                if c.school_class and c.school_class.grade_level_id
            )
            grade_levels = GradeLevel.objects.filter(id__in=gl_ids).order_by('number')
        except Exception:
            grade_levels = GradeLevel.objects.none()
    else:
        grade_levels = GradeLevel.objects.none()

    data = [{'id': gl.id, 'number': gl.number, 'name': str(gl)} for gl in grade_levels]
    return Response(data)


# ─── Аннотации учебника ────────────────────────────────────────────────────────

@api_view(['GET', 'PUT'])
@permission_classes([IsAuthenticated, PasswordChanged])
def textbook_annotations(request, session_id, slide_id):
    """
    GET  — все аннотации текущего студента на слайде (все страницы).
    PUT  — сохранить аннотацию страницы {page_number, strokes}.
    """
    session = get_object_or_404(LessonSession, id=session_id)
    slide   = get_object_or_404(Slide, id=slide_id)

    if request.method == 'GET':
        qs = TextbookAnnotation.objects.filter(session=session, slide=slide, student=request.user)
        return Response([{'page_number': a.page_number, 'strokes': a.strokes} for a in qs])

    page_number = request.data.get('page_number')
    if page_number is None:
        return Response({'error': 'page_number required'}, status=400)
    obj, _ = TextbookAnnotation.objects.update_or_create(
        session=session, slide=slide, student=request.user, page_number=int(page_number),
        defaults={'strokes': request.data.get('strokes', [])},
    )
    return Response({'page_number': obj.page_number, 'strokes': obj.strokes})


# ─── LessonAssignment views ────────────────────────────────────────────────────

@api_view(['GET', 'POST'])
@permission_classes([IsAuthenticated, PasswordChanged])
def lesson_assignments(request):
    """
    GET  — список заданий (учитель видит свои выдачи, ученик — выдачи ему)
    POST — создать выдачу (только учитель/admin)
    """
    if request.method == 'GET':
        user = request.user
        if _is_staff(user):
            qs = LessonAssignment.objects.filter(assigned_by=user).select_related(
                'lesson', 'school_class', 'student', 'assigned_by'
            )
        else:
            # Ученик: выдачи на его класс ИЛИ лично на него
            try:
                student_class = user.student_profile.school_class
            except Exception:
                student_class = None
            qs = LessonAssignment.objects.filter(
                django_models.Q(student=user) |
                (django_models.Q(school_class=student_class) if student_class else django_models.Q(pk__in=[]))
            ).select_related('lesson', 'school_class', 'student', 'assigned_by')
        return Response(LessonAssignmentSerializer(qs, many=True, context=_ctx(request)).data)

    # POST
    if not _is_staff(request.user):
        return Response({'error': 'Нет доступа'}, status=403)

    lesson_id = request.data.get('lesson')
    school_class_id = request.data.get('school_class')
    student_id = request.data.get('student')
    due_date = request.data.get('due_date') or None

    if not lesson_id:
        return Response({'error': 'lesson required'}, status=400)
    if not school_class_id and not student_id:
        return Response({'error': 'school_class or student required'}, status=400)

    lesson = get_object_or_404(Lesson, id=lesson_id)
    assignment = LessonAssignment.objects.create(
        lesson=lesson,
        school_class_id=school_class_id or None,
        student_id=student_id or None,
        assigned_by=request.user,
        due_date=due_date,
    )
    return Response(LessonAssignmentSerializer(assignment, context=_ctx(request)).data, status=201)


@api_view(['DELETE'])
@permission_classes([IsAuthenticated, PasswordChanged])
def lesson_assignment_detail(request, assignment_id):
    """DELETE — отозвать выдачу (только создавший или admin)."""
    assignment = get_object_or_404(LessonAssignment, id=assignment_id)
    if assignment.assigned_by != request.user and not request.user.is_admin:
        return Response({'error': 'Нет доступа'}, status=403)
    assignment.delete()
    return Response(status=204)
